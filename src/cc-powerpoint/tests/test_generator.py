"""Tests for the PowerPoint generator."""

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation

from src.parser import parse_markdown
from src.pptx_generator import generate_pptx
from src.themes import get_theme


def _generate_and_open(markdown: str, theme_name: str = "paper") -> Presentation:
    """Helper: parse markdown, generate pptx, return opened Presentation."""
    slides = parse_markdown(markdown)
    theme = get_theme(theme_name)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    generate_pptx(slides, theme, output_path)
    prs = Presentation(str(output_path))
    output_path.unlink(missing_ok=True)
    return prs


class TestBasicGeneration:
    """Tests for basic PPTX generation."""

    def test_creates_file(self):
        md = "---\n# Hello World\n## Test\n---"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = Path(f.name)

        slides = parse_markdown(md)
        theme = get_theme("paper")
        generate_pptx(slides, theme, output_path)

        assert output_path.exists()
        assert output_path.stat().st_size > 0
        output_path.unlink(missing_ok=True)

    def test_slide_count(self):
        md = "---\n# Title\n---\n# Slide 2\n\n- A\n- B\n---\n# Slide 3\n---"
        prs = _generate_and_open(md)
        assert len(prs.slides) == 3

    def test_widescreen_dimensions(self):
        md = "---\n# Title\n---"
        prs = _generate_and_open(md)
        # 16:9 = 13.333" x 7.5"
        assert prs.slide_width > prs.slide_height


class TestThemes:
    """Tests for theme application."""

    def test_all_themes_generate(self):
        md = "---\n# Title\n## Subtitle\n---\n# Content\n\n- Bullet\n---"
        for theme_name in ["boardroom", "paper", "terminal", "spark",
                           "thesis", "obsidian", "blueprint"]:
            prs = _generate_and_open(md, theme_name)
            assert len(prs.slides) == 2, f"Failed for theme: {theme_name}"

    def test_invalid_theme_raises(self):
        try:
            get_theme("nonexistent")
            assert False, "Should have raised ValueError"
        except ValueError:
            pass


class TestLayoutGeneration:
    """Tests for specific layout types producing valid output."""

    def test_title_slide_has_shapes(self):
        md = "---\n# Presentation\n## Author\n---"
        prs = _generate_and_open(md)
        slide = prs.slides[0]
        # Title slide should have text boxes and accent line
        assert len(slide.shapes) >= 2

    def test_content_slide_has_text(self):
        md = "---\n# Title\n---\n# Points\n\n- Alpha\n- Beta\n---"
        prs = _generate_and_open(md)
        slide = prs.slides[1]
        # Should have heading + bullet text box + accent line
        assert len(slide.shapes) >= 2

    def test_table_slide(self):
        md = "---\n# T\n---\n# Data\n\n| X | Y |\n|---|---|\n| 1 | 2 |\n---"
        prs = _generate_and_open(md)
        slide = prs.slides[1]
        has_table = any(shape.has_table for shape in slide.shapes)
        assert has_table

    def test_code_slide_has_shapes(self):
        md = "---\n# T\n---\n# Code\n\n```python\nx = 1\n```\n---"
        prs = _generate_and_open(md)
        slide = prs.slides[1]
        # Should have heading, accent line, code bg, code text, lang label
        assert len(slide.shapes) >= 3

    def test_speaker_notes(self):
        md = "---\n# T\n---\n# Slide\n\n- Point\n\n> My notes here\n---"
        prs = _generate_and_open(md)
        slide = prs.slides[1]
        notes = slide.notes_slide.notes_text_frame.text
        assert "My notes here" in notes
