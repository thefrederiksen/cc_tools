"""PowerPoint generator using python-pptx.

Converts parsed Slide dataclasses into a .pptx file with themed formatting.
Uses blank layouts with programmatic text boxes and shapes for reliability.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .parser import Slide, SlideLayout, TableData, ImageData
from .themes import PresentationTheme


# Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Standard margins
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _set_slide_background(pptx_slide, hex_color: str) -> None:
    """Set the background color of a slide."""
    background = pptx_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_text_box(
    pptx_slide,
    left,
    top,
    width,
    height,
    text: str,
    font_name: str,
    font_size: int,
    font_color: str,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> None:
    """Add a text box to a slide with specified formatting."""
    tx_box = pptx_slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = _hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = alignment


def _add_speaker_notes(pptx_slide, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    if notes_text:
        notes_slide = pptx_slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def generate_pptx(
    slides: list[Slide],
    theme: PresentationTheme,
    output_path: Path,
    input_dir: Path | None = None,
) -> None:
    """Generate a PowerPoint file from parsed slides.

    Args:
        slides: List of parsed Slide dataclasses
        theme: Presentation theme to apply
        output_path: Path for the output .pptx file
        input_dir: Directory of the input markdown file (for resolving image paths)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use the blank layout (index 6 in default template)
    blank_layout = prs.slide_layouts[6]

    for slide_data in slides:
        pptx_slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(pptx_slide, theme.colors.background)

        # Build the slide based on layout type
        builder = _LAYOUT_BUILDERS.get(slide_data.layout, _build_blank)
        builder(pptx_slide, slide_data, theme, input_dir)

        # Add speaker notes
        _add_speaker_notes(pptx_slide, slide_data.speaker_notes)

    prs.save(str(output_path))


# -- Layout Builders --

def _build_title_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a title slide with centered title and subtitle."""
    # Title - centered vertically in upper portion
    title_top = Inches(2.2)
    title_height = Inches(1.5)
    _add_text_box(
        pptx_slide,
        left=MARGIN_LEFT,
        top=title_top,
        width=CONTENT_WIDTH,
        height=title_height,
        text=slide.title,
        font_name=theme.fonts.heading,
        font_size=theme.title_font_size,
        font_color=theme.colors.heading,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if slide.subtitle:
        subtitle_top = title_top + title_height + Inches(0.2)
        _add_text_box(
            pptx_slide,
            left=MARGIN_LEFT,
            top=subtitle_top,
            width=CONTENT_WIDTH,
            height=Inches(0.8),
            text=slide.subtitle,
            font_name=theme.fonts.body,
            font_size=theme.subtitle_font_size,
            font_color=theme.colors.text,
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )

    # Accent line under title
    line_top = Inches(2.0)
    line_width = Inches(3.0)
    line_left = (SLIDE_WIDTH - line_width) // 2
    shape = pptx_slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        line_left, line_top, line_width, Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(theme.colors.accent)
    shape.line.fill.background()


def _build_section_header(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a section header slide with centered title."""
    _add_text_box(
        pptx_slide,
        left=MARGIN_LEFT,
        top=Inches(2.8),
        width=CONTENT_WIDTH,
        height=Inches(1.5),
        text=slide.title,
        font_name=theme.fonts.heading,
        font_size=theme.heading_font_size + 4,
        font_color=theme.colors.heading,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Accent line
    line_width = Inches(2.0)
    line_left = (SLIDE_WIDTH - line_width) // 2
    shape = pptx_slide.shapes.add_shape(
        1,
        line_left, Inches(4.4), line_width, Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(theme.colors.accent)
    shape.line.fill.background()


def _build_content_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a title + bullet content slide."""
    # Title at top
    _add_slide_heading(pptx_slide, slide.title, theme)

    # Bullets
    content_top = Inches(1.8)
    content_height = SLIDE_HEIGHT - content_top - Inches(0.6)

    tx_box = pptx_slide.shapes.add_textbox(
        MARGIN_LEFT, content_top, CONTENT_WIDTH, content_height,
    )
    tf = tx_box.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(slide.bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = theme.fonts.body
        p.font.size = Pt(theme.body_font_size)
        p.font.color.rgb = _hex_to_rgb(theme.colors.text)
        p.space_after = Pt(8)
        p.level = 0

        # Sub-bullets
        if idx in slide.sub_bullets:
            for sub in slide.sub_bullets[idx]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.name = theme.fonts.body
                sp.font.size = Pt(theme.body_font_size - 2)
                sp.font.color.rgb = _hex_to_rgb(theme.colors.text)
                sp.space_after = Pt(4)
                sp.level = 1


def _build_table_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a title + table slide."""
    _add_slide_heading(pptx_slide, slide.title, theme)

    table_data = slide.table
    if not table_data:
        return

    rows = len(table_data.rows) + (1 if table_data.headers else 0)
    cols = len(table_data.headers) if table_data.headers else (
        len(table_data.rows[0]) if table_data.rows else 1
    )

    table_top = Inches(1.8)
    table_height = min(Inches(0.4) * rows, SLIDE_HEIGHT - table_top - Inches(0.8))

    tbl_shape = pptx_slide.shapes.add_table(
        rows, cols, MARGIN_LEFT, table_top, CONTENT_WIDTH, table_height,
    )
    tbl = tbl_shape.table

    # Style header row
    row_idx = 0
    if table_data.headers:
        for col_idx, header in enumerate(table_data.headers):
            cell = tbl.cell(0, col_idx)
            cell.text = header
            _style_table_cell(cell, theme, is_header=True)
        row_idx = 1

    # Data rows
    for data_row in table_data.rows:
        for col_idx, value in enumerate(data_row):
            if col_idx < cols:
                cell = tbl.cell(row_idx, col_idx)
                cell.text = value
                _style_table_cell(cell, theme, is_header=False)
        row_idx += 1


def _style_table_cell(cell, theme: PresentationTheme, is_header: bool) -> None:
    """Apply theme styling to a table cell."""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = theme.fonts.body
        paragraph.font.size = Pt(theme.body_font_size - 2)
        if is_header:
            paragraph.font.bold = True
            paragraph.font.color.rgb = _hex_to_rgb(theme.colors.background)
        else:
            paragraph.font.color.rgb = _hex_to_rgb(theme.colors.text)

    # Cell background
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(theme.colors.primary)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _build_code_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a title + code block slide."""
    _add_slide_heading(pptx_slide, slide.title, theme)

    code_top = Inches(1.8)
    code_height = SLIDE_HEIGHT - code_top - Inches(0.8)
    code_left = MARGIN_LEFT
    code_width = CONTENT_WIDTH

    # Code background rectangle
    bg_shape = pptx_slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        code_left, code_top, code_width, code_height,
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _hex_to_rgb(theme.colors.code_bg)
    bg_shape.line.fill.background()

    # Language label in top-right
    if slide.code_language:
        _add_text_box(
            pptx_slide,
            left=code_left + code_width - Inches(2.0),
            top=code_top + Inches(0.1),
            width=Inches(1.8),
            height=Inches(0.3),
            text=slide.code_language,
            font_name=theme.fonts.code,
            font_size=10,
            font_color=theme.colors.accent,
            alignment=PP_ALIGN.RIGHT,
        )

    # Code text
    padding = Inches(0.3)
    tx_box = pptx_slide.shapes.add_textbox(
        code_left + padding,
        code_top + Inches(0.4),
        code_width - padding * 2,
        code_height - Inches(0.5),
    )
    tf = tx_box.text_frame
    tf.word_wrap = True

    # Split code into lines and add each as a paragraph
    lines = slide.code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = theme.fonts.code
        p.font.size = Pt(theme.code_font_size)
        p.font.color.rgb = _hex_to_rgb(theme.colors.text)
        p.space_after = Pt(2)


def _build_image_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a title + image slide."""
    _add_slide_heading(pptx_slide, slide.title, theme)
    _place_image(pptx_slide, slide.image, input_dir, top=Inches(1.8))


def _build_blank_image_slide(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a full-image slide with no heading."""
    _place_image(pptx_slide, slide.image, input_dir, top=Inches(0.5))


def _build_blank(
    pptx_slide,
    slide: Slide,
    theme: PresentationTheme,
    input_dir: Path | None,
) -> None:
    """Build a blank/fallback slide - just show raw text if any."""
    if slide.raw_markdown:
        _add_text_box(
            pptx_slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=SLIDE_HEIGHT - MARGIN_TOP * 2,
            text=slide.raw_markdown,
            font_name=theme.fonts.body,
            font_size=theme.body_font_size,
            font_color=theme.colors.text,
        )


# -- Shared Helpers --

def _add_slide_heading(pptx_slide, title: str, theme: PresentationTheme) -> None:
    """Add a standard heading bar at the top of a content slide."""
    # Title text
    _add_text_box(
        pptx_slide,
        left=MARGIN_LEFT,
        top=Inches(0.5),
        width=CONTENT_WIDTH,
        height=Inches(0.8),
        text=title,
        font_name=theme.fonts.heading,
        font_size=theme.heading_font_size,
        font_color=theme.colors.heading,
        bold=True,
    )

    # Accent underline
    shape = pptx_slide.shapes.add_shape(
        1,
        MARGIN_LEFT, Inches(1.35), Inches(4.0), Inches(0.03),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(theme.colors.accent)
    shape.line.fill.background()


def _place_image(
    pptx_slide,
    image: ImageData | None,
    input_dir: Path | None,
    top=Inches(1.8),
) -> None:
    """Place an image on the slide, centered and scaled to fit."""
    if not image:
        return

    # Resolve image path
    img_path = Path(image.src)
    if not img_path.is_absolute() and input_dir:
        img_path = input_dir / img_path

    if not img_path.exists():
        # Show placeholder text if image not found
        _add_text_box(
            pptx_slide,
            left=MARGIN_LEFT,
            top=top,
            width=CONTENT_WIDTH,
            height=Inches(1.0),
            text=f"[Image not found: {image.src}]",
            font_name="Segoe UI",
            font_size=14,
            font_color="#999999",
            alignment=PP_ALIGN.CENTER,
        )
        return

    # Available space
    max_width = CONTENT_WIDTH - Inches(1.0)
    max_height = SLIDE_HEIGHT - top - Inches(0.8)

    # Add the image - python-pptx will use its native size
    pic = pptx_slide.shapes.add_picture(
        str(img_path), MARGIN_LEFT, top,
    )

    # Scale to fit within available space
    img_width = pic.width
    img_height = pic.height

    width_ratio = max_width / img_width if img_width > max_width else 1.0
    height_ratio = max_height / img_height if img_height > max_height else 1.0
    scale = min(width_ratio, height_ratio)

    if scale < 1.0:
        pic.width = int(img_width * scale)
        pic.height = int(img_height * scale)

    # Center horizontally
    pic.left = int(MARGIN_LEFT + (CONTENT_WIDTH - pic.width) / 2)


# Layout builder dispatch table
_LAYOUT_BUILDERS = {
    SlideLayout.TITLE: _build_title_slide,
    SlideLayout.SECTION_HEADER: _build_section_header,
    SlideLayout.TITLE_CONTENT: _build_content_slide,
    SlideLayout.TITLE_TABLE: _build_table_slide,
    SlideLayout.TITLE_CODE: _build_code_slide,
    SlideLayout.TITLE_IMAGE: _build_image_slide,
    SlideLayout.BLANK_IMAGE: _build_blank_image_slide,
    SlideLayout.BLANK: _build_blank,
}
